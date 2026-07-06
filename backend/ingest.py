"""Primary implementation for the raw-source ingestion pipeline.

This module contains the actual document ingestion workflow. The
backend/process_raw_sources.py entry point is a thin compatibility wrapper
that calls this implementation.
"""
from __future__ import annotations

import argparse
import csv
import json
from pathlib import Path
from typing import Iterator, List

from bs4 import BeautifulSoup
from docx import Document as DocxDocument
from openpyxl import load_workbook
from pptx import Presentation as PptPresentation
from pypdf import PdfReader


SUPPORTED_EXTENSIONS = {
    ".txt",
    ".md",
    ".html",
    ".htm",
    ".xml",
    ".json",
    ".csv",
    ".pdf",
    ".docx",
    ".pptx",
    ".xlsx",
}

DEFAULT_OUTPUT_ROOT = Path(__file__).resolve().parent / "json"


def iter_supported_files(root: Path) -> Iterator[Path]:
    if not root.exists():
        raise FileNotFoundError(f"Path does not exist: {root}")

    if root.is_file():
        if root.suffix.lower() in SUPPORTED_EXTENSIONS:
            yield root
        return

    for path in sorted(root.rglob("*")):
        if path.is_file() and path.suffix.lower() in SUPPORTED_EXTENSIONS:
            yield path


def build_output_path(source: Path, output_root: Path, input_path: Path) -> Path:
    if source.is_file():
        return output_root / f"{source.stem}.json"

    relative_path = input_path.relative_to(source)
    return output_root / relative_path.with_suffix(".json")


def extract_text(path: Path) -> str:
    suffix = path.suffix.lower()

    if suffix in {".txt", ".md"}:
        return path.read_text(encoding="utf-8", errors="ignore")

    if suffix in {".html", ".htm"}:
        html = path.read_text(encoding="utf-8", errors="ignore")
        soup = BeautifulSoup(html, "html.parser")
        return soup.get_text("\n", strip=True)

    if suffix == ".xml":
        xml = path.read_text(encoding="utf-8", errors="ignore")
        soup = BeautifulSoup(xml, "xml")
        return soup.get_text("\n", strip=True)

    if suffix == ".json":
        data = json.loads(path.read_text(encoding="utf-8", errors="ignore"))
        return json.dumps(data, indent=2, ensure_ascii=False)

    if suffix == ".csv":
        with path.open("r", encoding="utf-8-sig", errors="ignore", newline="") as handle:
            rows = list(csv.reader(handle))
        return "\n".join(",".join(row) for row in rows)

    if suffix == ".pdf":
        reader = PdfReader(str(path))
        pages = [page.extract_text() or "" for page in reader.pages]
        return "\n".join(page for page in pages if page)

    if suffix == ".docx":
        document = DocxDocument(str(path))
        paragraphs = [paragraph.text for paragraph in document.paragraphs if paragraph.text]
        return "\n".join(paragraphs)

    if suffix == ".pptx":
        presentation = PptPresentation(str(path))
        slides: List[str] = []
        for slide in presentation.slides:
            text_fragments = []
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text_fragments.append(shape.text)
            if text_fragments:
                slides.append("\n".join(text_fragments))
        return "\n\n".join(slides)

    if suffix == ".xlsx":
        workbook = load_workbook(filename=str(path), read_only=True, data_only=True)
        rows: List[str] = []
        for row in workbook.active.iter_rows(values_only=True):
            cleaned = [str(cell).strip() if cell is not None else "" for cell in row]
            rows.append("\t".join(cleaned))
        return "\n".join(rows)

    raise ValueError(f"Unsupported file type: {path}")


def ingest(
    source: Path,
    *,
    output_root: Path | None = None,
    output_file: Path | None = None,
) -> List[dict]:
    documents: List[dict] = []

    for path in iter_supported_files(source):
        document = {
            "path": str(path),
            "text": extract_text(path),
        }
        documents.append(document)

        if output_root:
            output_path = build_output_path(source, output_root, path)
            output_path.parent.mkdir(parents=True, exist_ok=True)
            output_path.write_text(
                json.dumps(document, indent=2, ensure_ascii=False),
                encoding="utf-8",
            )

    if output_file:
        output_file.parent.mkdir(parents=True, exist_ok=True)
        output_file.write_text(
            json.dumps(documents, indent=2, ensure_ascii=False),
            encoding="utf-8",
        )

    return documents


def resolve_output_path(path: Path | None) -> Path | None:
    if path is None:
        return None
    if path.is_absolute():
        return path

    repo_root = Path(__file__).resolve().parent.parent
    return repo_root / path


def main() -> None:
    parser = argparse.ArgumentParser(description="Ingest supported documents into text.")
    parser.add_argument("source", type=Path, help="File or directory to ingest")
    parser.add_argument(
        "--output-root",
        type=Path,
        default=DEFAULT_OUTPUT_ROOT,
        help="Directory for per-document JSON output (default: ./backend/json)",
    )
    parser.add_argument(
        "--output",
        "-o",
        type=Path,
        help="Optional JSON file that collects all ingested documents",
    )
    args = parser.parse_args()

    args.output_root = resolve_output_path(args.output_root)
    args.output = resolve_output_path(args.output)

    if args.output:
        documents = ingest(args.source, output_root=None, output_file=args.output)
    else:
        documents = ingest(args.source, output_root=args.output_root)

    if not args.output:
        for document in documents:
            print(f"=== {document['path']} ===")
            print(document["text"])
            print()


if __name__ == "__main__":
    main()